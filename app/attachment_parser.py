"""
v1.6.0+: Парсеры файлов-attachments из КБ Респект.Чата.

Извлекают текст из скачанных файлов для индексации в FTS.

Поддерживаемые форматы (parseable):
    pdf, docx, txt, md, html, rtf, pptx, xlsx, xls, csv

Не парсятся (только метаданные/ссылка):
    video, audio, image, doc, ppt (старые бинарные офисные форматы),
    другие неизвестные

Для старого MS Office (.doc, .ppt) понадобились бы внешние утилиты
(antiword/catppt) или LibreOffice — не используем, чтобы не раздувать
зависимости. Файлы в этих форматах остаются как ссылки в ответе Leo.
"""
from __future__ import annotations

import csv
import io
import logging
from typing import Optional


log = logging.getLogger("attachment_parser")


# Форматы которые мы умеем парсить (lowercase)
PARSEABLE_TYPES = {"pdf", "docx", "txt", "md", "html", "rtf", "pptx", "xlsx", "xls", "csv"}


def get_file_type(name: str, declared_type: Optional[str] = None) -> str:
    """Определить тип файла по расширению или declared_type."""
    if declared_type:
        t = declared_type.strip().lower()
        t = {
            "jpeg": "image", "jpg": "image", "png": "image", "gif": "image", "webp": "image",
            "mp4": "video", "mov": "video", "avi": "video", "mkv": "video", "webm": "video",
            "mp3": "audio", "wav": "audio", "ogg": "audio", "m4a": "audio",
            "document": "docx",
            "text": "txt", "plain": "txt",
            "markdown": "md",
            "spreadsheet": "xlsx",
            "presentation": "pptx",
        }.get(t, t)
        if t in PARSEABLE_TYPES or t in {"video", "audio", "image", "doc", "ppt", "other"}:
            return t

    if not name:
        return "other"

    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    return {
        # Парсим
        "pdf":  "pdf",
        "docx": "docx",
        "txt":  "txt",
        "log":  "txt",
        "md":   "md",
        "markdown": "md",
        "html": "html",
        "htm":  "html",
        "rtf":  "rtf",
        "pptx": "pptx",
        "xlsx": "xlsx",
        "xlsm": "xlsx",
        "xls":  "xls",
        "csv":  "csv",
        "tsv":  "csv",
        # Не парсим — старые бинарные форматы офиса
        "doc":  "doc",
        "ppt":  "ppt",
        # Медиа
        "jpg":  "image", "jpeg": "image", "png": "image", "gif": "image", "webp": "image", "svg": "image", "bmp": "image",
        "mp4":  "video", "mov":  "video", "avi": "video", "mkv":  "video", "webm": "video",
        "mp3":  "audio", "wav":  "audio", "ogg": "audio", "m4a":  "audio",
        # Прочее
        "json": "other",
        "zip":  "other",
        "rar":  "other",
        "7z":   "other",
    }.get(ext, "other")


def is_parseable(file_type: str) -> bool:
    return file_type.lower() in PARSEABLE_TYPES


# -----------------------------------------------------------------------------
# Парсеры
# -----------------------------------------------------------------------------
def _parse_pdf(blob: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(blob))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            log.warning("pdf page %d extract failed: %s", i, e)
            txt = ""
        if txt.strip():
            parts.append(txt)
    return "\n\n".join(parts).strip()


def _parse_docx(blob: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(blob))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _parse_txt(blob: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return blob.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return blob.decode("utf-8", errors="replace").strip()


def _parse_md(blob: bytes) -> str:
    return _parse_txt(blob)


def _parse_html(blob: bytes) -> str:
    from app.html_utils import html_to_plain
    text = _parse_txt(blob)
    return html_to_plain(text)


def _parse_rtf(blob: bytes) -> str:
    """RTF → plain text через striprtf (pure Python)."""
    from striprtf.striprtf import rtf_to_text
    text = _parse_txt(blob)  # rtf — это текст в ASCII/cp1252/utf-8
    return rtf_to_text(text, errors="ignore").strip()


def _parse_pptx(blob: bytes) -> str:
    """PPTX → текст: слайды + заметки + таблицы."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            # Текстовые блоки
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_parts.append(txt)
            # Таблицы
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
        # Заметки докладчика
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes] {notes}")
        if slide_parts:
            parts.append(f"[Slide {i}]")
            parts.extend(slide_parts)
    return "\n".join(parts).strip()


def _parse_xlsx(blob: bytes) -> str:
    """XLSX → текст всех ячеек по листам."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(blob), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def _parse_xls(blob: bytes) -> str:
    """XLS (legacy) → текст через xlrd."""
    import xlrd  # type: ignore
    book = xlrd.open_workbook(file_contents=blob)
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"[Sheet: {sheet.name}]")
        for row_idx in range(sheet.nrows):
            cells = []
            for col_idx in range(sheet.ncols):
                v = sheet.cell_value(row_idx, col_idx)
                s = str(v).strip() if v not in (None, "") else ""
                if s:
                    cells.append(s)
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _parse_csv(blob: bytes) -> str:
    """CSV/TSV → текст. Авто-определение разделителя."""
    text = _parse_txt(blob)
    # Sniff delimiter из первых 4 KB
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        # Fallback: запятая
        class _D:
            delimiter = ","
        dialect = _D()  # type: ignore

    parts: list[str] = []
    reader = csv.reader(io.StringIO(text), dialect)
    for row in reader:
        cells = [c.strip() for c in row if c and c.strip()]
        if cells:
            parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


_PARSERS = {
    "pdf":  _parse_pdf,
    "docx": _parse_docx,
    "txt":  _parse_txt,
    "md":   _parse_md,
    "html": _parse_html,
    "rtf":  _parse_rtf,
    "pptx": _parse_pptx,
    "xlsx": _parse_xlsx,
    "xls":  _parse_xls,
    "csv":  _parse_csv,
}


def parse_blob(blob: bytes, file_type: str) -> str:
    parser = _PARSERS.get(file_type.lower())
    if parser is None:
        raise ValueError(f"file_type {file_type!r} not parseable")
    return parser(blob) or ""
