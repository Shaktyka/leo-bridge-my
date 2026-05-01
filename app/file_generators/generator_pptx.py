"""PPTX генератор — JSON через python-pptx."""
import json
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

log = logging.getLogger(__name__)


def generate_pptx(out_path: Path, content_md: str | None = None,
                  content_json: str | None = None, title: str | None = None) -> None:
    """Создать .pptx из JSON структуры.

    Формат JSON:
    {
      "slides": [
        {"title": "Заголовок", "bullets": ["item1", "item2"]},
        {"title": "Заголовок 2", "body": "Длинный текст параграфом"}
      ]
    }
    """
    if not content_json:
        raise ValueError("content_json is required for pptx format")

    data = json.loads(content_json) if isinstance(content_json, str) else content_json
    slides = data.get("slides") or []
    if not slides:
        raise ValueError("pptx: slides list is empty")

    prs = Presentation()

    # Title slide если задан title
    if title:
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title + Content
    for s_data in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = s_data.get("title", "")

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        if "bullets" in s_data:
            for i, bullet in enumerate(s_data["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
        elif "body" in s_data:
            tf.text = s_data["body"]

    prs.save(str(out_path))
    log.info("pptx generated: %s (%d slides)", out_path, len(slides))
